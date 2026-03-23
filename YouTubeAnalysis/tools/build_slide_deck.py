"""
build_slide_deck.py
Generates a professional .pptx slide deck from the analysis JSON.

Usage:
    python tools/build_slide_deck.py
    python tools/build_slide_deck.py --input .tmp/analysis_2026-03-23.json
    python tools/build_slide_deck.py --date 2026-03-23

Reads:
    .tmp/analysis_YYYY-MM-DD.json  (from analyze_trends.py, with agent-written narrative)
    config.json                     (email_subject_template for deck title)

Output (stdout): JSON summary
Writes:
    .tmp/charts/*.png               (intermediate chart images)
    .tmp/decks/youtube_trends_YYYY-MM-DD.pptx

Slides:
    1. Cover + Executive Summary
    2. Top 10 Trending Videos (table)
    3. View Velocity Chart
    4. Engagement Rate Chart
    5. Top Keywords Bar Chart
    6. Top Channels Table + Spotlights
    7. Transcript Themes
    8. Content Opportunity Gaps + Recommendations
"""

import argparse
import json
import sys
from datetime import date
from io import BytesIO
from pathlib import Path

PROJECT_ROOT = Path(__file__).parent.parent
TMP_DIR = PROJECT_ROOT / ".tmp"
CHARTS_DIR = TMP_DIR / "charts"
DECKS_DIR = TMP_DIR / "decks"

try:
    from dotenv import load_dotenv
except ImportError:
    print("ERROR: python-dotenv not installed. Run: pip install python-dotenv", file=sys.stderr)
    sys.exit(1)

load_dotenv(PROJECT_ROOT / ".env")

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
except ImportError:
    print("ERROR: matplotlib not installed. Run: pip install matplotlib", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------

DARK_BG = RGBColor(0x0F, 0x0F, 0x23)       # near-black navy
ACCENT = RGBColor(0xFF, 0x45, 0x00)          # YouTube-red orange
ACCENT2 = RGBColor(0x00, 0xC2, 0xA8)         # teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
MID_GREY = RGBColor(0x98, 0xA2, 0xB3)
DARK_TEXT = RGBColor(0x1D, 0x29, 0x39)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def load_config() -> dict:
    config_path = PROJECT_ROOT / "config.json"
    if not config_path.exists():
        return {}
    return json.loads(config_path.read_text(encoding="utf-8"))


def find_analysis(target_date: str | None) -> Path:
    if target_date:
        p = TMP_DIR / f"analysis_{target_date}.json"
        if not p.exists():
            print(f"ERROR: {p} not found", file=sys.stderr)
            sys.exit(1)
        return p
    candidates = sorted(TMP_DIR.glob("analysis_*.json"), reverse=True)
    if not candidates:
        print("ERROR: No analysis_*.json found in .tmp/ — run analyze_trends.py first", file=sys.stderr)
        sys.exit(1)
    return candidates[0]


def fmt_number(n: int | float) -> str:
    """Format large numbers: 1234567 -> '1.2M', 12345 -> '12.3K'"""
    n = int(n)
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n / 1_000:.1f}K"
    return str(n)


def truncate(s: str, max_len: int) -> str:
    return s[:max_len] + "…" if len(s) > max_len else s


def blank_slide(prs: Presentation, layout_index: int = 6) -> object:
    """Add a slide with the blank layout."""
    return prs.slides.add_slide(prs.slide_layouts[layout_index])


def set_slide_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def chart_to_image(fig) -> BytesIO:
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


# ---------------------------------------------------------------------------
# Chart generators
# ---------------------------------------------------------------------------

def make_hbar_chart(labels, values, title, color="#FF4500", value_fmt=None, figsize=(9, 5)):
    """Generic horizontal bar chart. Returns a BytesIO PNG."""
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor("#0F0F23")
    ax.set_facecolor("#1A1A35")

    y_pos = range(len(labels))
    bars = ax.barh(list(y_pos), values, color=color, height=0.6, zorder=3)

    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(labels, color="white", fontsize=8)
    ax.invert_yaxis()
    ax.set_xlabel("", color="white")
    ax.set_title(title, color="white", fontsize=11, pad=10)
    ax.tick_params(colors="white", labelsize=8)
    ax.xaxis.label.set_color("white")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#333355")
    ax.spines["bottom"].set_color("#333355")
    ax.tick_params(axis="x", colors="#888899")
    ax.grid(axis="x", color="#333355", linestyle="--", alpha=0.5, zorder=0)

    # Value labels on bars
    for bar, val in zip(bars, values):
        label = value_fmt(val) if value_fmt else str(val)
        ax.text(
            bar.get_width() + max(values) * 0.01,
            bar.get_y() + bar.get_height() / 2,
            label,
            va="center", ha="left", color="white", fontsize=7,
        )

    plt.tight_layout()
    return chart_to_image(fig)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_slide_cover(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)

    # Accent bar left edge
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)

    # Title
    add_textbox(slide, Inches(0.4), Inches(0.5), Inches(8), Inches(1.2),
                "AI & Automation", font_size=40, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(1.6), Inches(8), Inches(0.7),
                "YouTube Trend Report", font_size=28, bold=False, color=ACCENT)
    add_textbox(slide, Inches(0.4), Inches(2.2), Inches(6), Inches(0.5),
                f"Week of {analysis.get('generated_date', date.today().isoformat())}",
                font_size=14, color=MID_GREY)

    # Stat callout boxes
    stats = [
        (fmt_number(analysis.get("total_videos_analyzed", 0)), "Videos Analyzed"),
        (f"{analysis.get('period_days', 14)} days", "Data Period"),
        (fmt_number(analysis.get("fastest_growing_video", {}).get("view_velocity", 0)) + "/day",
         "Top View Velocity"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(1.0)
    box_top = Inches(3.2)
    for i, (stat_val, stat_label) in enumerate(stats):
        left = Inches(0.4) + i * (box_w + Inches(0.3))
        add_rect(slide, left, box_top, box_w, box_h, RGBColor(0x1A, 0x1A, 0x35))
        add_textbox(slide, left + Inches(0.1), box_top + Inches(0.05),
                    box_w - Inches(0.2), Inches(0.5),
                    stat_val, font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), box_top + Inches(0.55),
                    box_w - Inches(0.2), Inches(0.35),
                    stat_label, font_size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Executive summary
    summary = analysis.get("executive_summary", "")
    if summary:
        add_textbox(slide, Inches(0.4), Inches(4.5), Inches(9), Inches(2.5),
                    summary, font_size=11, color=LIGHT_GREY)

    # Decorative accent line
    add_rect(slide, Inches(0.4), Inches(4.35), Inches(9), Inches(0.03), ACCENT)


def build_slide_top_videos(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                "Top 10 Trending Videos", font_size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.35),
                f"Sorted by view count • Past {analysis.get('period_days', 14)} days",
                font_size=11, color=MID_GREY)

    top_videos = analysis.get("top_videos", [])[:10]

    # Table headers
    cols = ["#", "Title", "Channel", "Views", "Engagement", "Published"]
    col_widths = [Inches(0.3), Inches(4.5), Inches(2.2), Inches(1.1), Inches(1.3), Inches(1.2)]
    row_h = Inches(0.52)
    header_top = Inches(1.1)
    left_start = Inches(0.3)

    # Header row background
    add_rect(slide, left_start, header_top, sum(col_widths), row_h, ACCENT)

    x = left_start
    for col, w in zip(cols, col_widths):
        add_textbox(slide, x + Inches(0.05), header_top + Inches(0.1),
                    w - Inches(0.1), row_h - Inches(0.1),
                    col, font_size=9, bold=True, color=WHITE)
        x += w

    # Data rows
    for row_idx, video in enumerate(top_videos):
        row_top = header_top + (row_idx + 1) * row_h
        row_color = RGBColor(0x1A, 0x1A, 0x35) if row_idx % 2 == 0 else RGBColor(0x14, 0x14, 0x2A)
        add_rect(slide, left_start, row_top, sum(col_widths), row_h, row_color)

        eng = video.get("engagement_rate", 0)
        eng_color = (RGBColor(0x12, 0xB7, 0x6A) if eng >= 3
                     else RGBColor(0xFF, 0xAB, 0x00) if eng >= 1
                     else RGBColor(0xF0, 0x44, 0x38))

        cells = [
            (str(row_idx + 1), WHITE),
            (truncate(video.get("title", ""), 55), LIGHT_GREY),
            (truncate(video.get("channel_title", ""), 28), MID_GREY),
            (fmt_number(video.get("view_count", 0)), WHITE),
            (f"{eng:.1f}%", eng_color),
            (video.get("published_at", "")[:10], MID_GREY),
        ]
        x = left_start
        for (cell_text, cell_color), w in zip(cells, col_widths):
            add_textbox(slide, x + Inches(0.05), row_top + Inches(0.1),
                        w - Inches(0.1), row_h - Inches(0.15),
                        cell_text, font_size=8, color=cell_color)
            x += w

    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(10), Inches(0.3),
                "Engagement = (Likes + Comments) / Views × 100. Green >3%, Amber 1-3%, Red <1%",
                font_size=8, color=MID_GREY)


def build_slide_view_velocity(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT2)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                "View Velocity Ranking", font_size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.35),
                "Views ÷ Days Since Publish — normalises for age, reveals what's actually growing now",
                font_size=11, color=MID_GREY)

    ranking = analysis.get("view_velocity_ranking", [])[:12]
    if not ranking:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                    "No velocity data available.", font_size=14, color=MID_GREY)
        return

    labels = [truncate(f"{v.get('channel_title', '')} — {v.get('title', '')}", 60) for v in ranking]
    values = [v.get("view_velocity", 0) for v in ranking]

    buf = make_hbar_chart(labels, values, "", color="#FF4500",
                          value_fmt=lambda v: f"{fmt_number(v)}/d", figsize=(10, 5))
    slide.shapes.add_picture(buf, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8))


def build_slide_engagement(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT2)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                "Engagement Rate Ranking", font_size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.35),
                f"Top videos by engagement rate (min {analysis.get('total_videos_analyzed', 0)} views threshold). "
                "High engagement signals algorithmic boost.",
                font_size=11, color=MID_GREY)

    ranking = analysis.get("engagement_ranking", [])[:12]
    if not ranking:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                    "No engagement data available.", font_size=14, color=MID_GREY)
        return

    labels = [truncate(v.get("title", ""), 60) for v in ranking]
    values = [v.get("engagement_rate", 0) for v in ranking]

    buf = make_hbar_chart(labels, values, "", color="#00C2A8",
                          value_fmt=lambda v: f"{v:.1f}%", figsize=(10, 5))
    slide.shapes.add_picture(buf, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8))


def build_slide_keywords(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, RGBColor(0x79, 0x5E, 0xC6))

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                "Top Keywords in Video Titles", font_size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.35),
                "Most frequent words across all discovered titles this period (stop words excluded)",
                font_size=11, color=MID_GREY)

    keyword_freq = analysis.get("keyword_frequency", {})
    if not keyword_freq:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                    "No keyword data available.", font_size=14, color=MID_GREY)
        return

    items = list(keyword_freq.items())[:20]
    labels = [k for k, _ in items]
    values = [v for _, v in items]
    median_val = sorted(values)[len(values) // 2]
    colors = ["#795EC6" if v > median_val else "#4A3A80" for v in values]

    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor("#0F0F23")
    ax.set_facecolor("#1A1A35")
    y_pos = range(len(labels))
    bars = ax.barh(list(y_pos), values, color=colors, height=0.6, zorder=3)
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(labels, color="white", fontsize=9)
    ax.invert_yaxis()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#333355")
    ax.spines["bottom"].set_color("#333355")
    ax.tick_params(colors="#888899", labelsize=8)
    ax.grid(axis="x", color="#333355", linestyle="--", alpha=0.5, zorder=0)
    for bar, val in zip(bars, values):
        ax.text(bar.get_width() + max(values) * 0.01, bar.get_y() + bar.get_height() / 2,
                str(val), va="center", ha="left", color="white", fontsize=8)

    legend_patches = [
        mpatches.Patch(color="#795EC6", label=f"Above median ({median_val} mentions)"),
        mpatches.Patch(color="#4A3A80", label="At or below median"),
    ]
    ax.legend(handles=legend_patches, loc="lower right", framealpha=0.3,
              labelcolor="white", fontsize=8)
    plt.tight_layout()
    buf = chart_to_image(fig)
    slide.shapes.add_picture(buf, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8))


def build_slide_top_channels(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                "Top Channels to Watch", font_size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.35),
                "Ranked by total views in dataset this period",
                font_size=11, color=MID_GREY)

    top_channels = analysis.get("top_channels", [])[:8]
    spotlights = analysis.get("channel_spotlights", [])

    # Table
    cols = ["Channel", "Subscribers", "Videos", "Avg Views/Video", "Top Video"]
    col_widths = [Inches(2.4), Inches(1.3), Inches(0.9), Inches(1.6), Inches(3.8)]
    row_h = Inches(0.52)
    header_top = Inches(1.1)
    left_start = Inches(0.3)

    add_rect(slide, left_start, header_top, sum(col_widths), row_h, ACCENT)
    x = left_start
    for col, w in zip(cols, col_widths):
        add_textbox(slide, x + Inches(0.05), header_top + Inches(0.1),
                    w - Inches(0.1), row_h - Inches(0.1),
                    col, font_size=9, bold=True, color=WHITE)
        x += w

    for row_idx, ch in enumerate(top_channels):
        row_top = header_top + (row_idx + 1) * row_h
        row_color = RGBColor(0x1A, 0x1A, 0x35) if row_idx % 2 == 0 else RGBColor(0x14, 0x14, 0x2A)
        add_rect(slide, left_start, row_top, sum(col_widths), row_h, row_color)
        cells = [
            truncate(ch.get("title", ""), 32),
            fmt_number(ch.get("subscriber_count", 0)),
            str(ch.get("videos_in_dataset", 0)),
            fmt_number(ch.get("avg_views_per_video", 0)),
            truncate(ch.get("top_video_title", ""), 52),
        ]
        x = left_start
        for cell_text, w in zip(cells, col_widths):
            add_textbox(slide, x + Inches(0.05), row_top + Inches(0.1),
                        w - Inches(0.1), row_h - Inches(0.15),
                        cell_text, font_size=8, color=LIGHT_GREY)
            x += w

    # Spotlights below table
    if spotlights:
        spot_top = header_top + (len(top_channels) + 1) * row_h + Inches(0.15)
        add_textbox(slide, left_start, spot_top, Inches(12), Inches(0.35),
                    "Channel Spotlights", font_size=11, bold=True, color=ACCENT)
        sp_box_w = Inches(4.1)
        for i, spotlight in enumerate(spotlights[:3]):
            sp_left = left_start + i * (sp_box_w + Inches(0.15))
            add_rect(slide, sp_left, spot_top + Inches(0.4), sp_box_w, Inches(0.75),
                     RGBColor(0x1A, 0x1A, 0x35))
            text = spotlight if isinstance(spotlight, str) else spotlight.get("text", str(spotlight))
            add_textbox(slide, sp_left + Inches(0.1), spot_top + Inches(0.45),
                        sp_box_w - Inches(0.2), Inches(0.65),
                        text, font_size=9, color=LIGHT_GREY)


def build_slide_transcript_themes(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT2)

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                "Transcript Theme Analysis", font_size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.65), Inches(10), Inches(0.35),
                "Topics creators are actually discussing — extracted from video transcripts (n-gram frequency)",
                font_size=11, color=MID_GREY)

    themes = analysis.get("transcript_themes", [])
    transcript_count = sum(
        1 for _ in analysis.get("top_videos", [])
        if analysis.get("total_videos_analyzed", 0) > 0
    )

    if not themes:
        add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1.5),
                    "No transcript data available for this run.\n"
                    "This may be because transcripts are disabled on most videos this period.",
                    font_size=14, color=MID_GREY)
        return

    # Left column: ranked list
    add_textbox(slide, Inches(0.4), Inches(1.1), Inches(5.5), Inches(0.35),
                "Top Discussed Topics", font_size=12, bold=True, color=WHITE)

    max_count = themes[0]["count"] if themes else 1

    for i, theme in enumerate(themes[:15]):
        top = Inches(1.5) + i * Inches(0.38)
        # Mini bar background
        add_rect(slide, Inches(0.4), top + Inches(0.05), Inches(5.5), Inches(0.28),
                 RGBColor(0x1A, 0x1A, 0x35))
        # Mini bar fill
        bar_width = max(Inches(0.1), Inches(5.5) * (theme["count"] / max_count))
        add_rect(slide, Inches(0.4), top + Inches(0.05), bar_width, Inches(0.28),
                 RGBColor(0x00, 0x80, 0x6A))
        # Label
        add_textbox(slide, Inches(0.5), top, Inches(4.5), Inches(0.35),
                    f"{i + 1}. {theme['phrase']}", font_size=9, color=WHITE)
        add_textbox(slide, Inches(5.3), top, Inches(0.6), Inches(0.35),
                    str(theme["count"]), font_size=8, color=MID_GREY)

    # Right column: context note
    add_textbox(slide, Inches(6.5), Inches(1.1), Inches(6.5), Inches(0.35),
                "What This Means", font_size=12, bold=True, color=WHITE)
    note = (
        "These phrases appear most frequently across all analyzed transcripts. "
        "They represent what creators are actually teaching and discussing — "
        "not just what they promise in titles.\n\n"
        "Cross-reference with the Keywords slide to find gaps between what's "
        "titled vs. what's being taught. Those gaps are your content opportunities."
    )
    add_textbox(slide, Inches(6.5), Inches(1.55), Inches(6.5), Inches(2.5),
                note, font_size=10, color=LIGHT_GREY)

    add_textbox(slide, Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
                f"Based on transcripts from top videos by view count this period",
                font_size=8, color=MID_GREY)


def build_slide_content_gaps(prs: Presentation, analysis: dict):
    slide = blank_slide(prs)
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, RGBColor(0x79, 0x5E, 0xC6))

    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                "Content Opportunity Gaps", font_size=24, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.35),
                "Topics discussed in transcripts but absent from video titles — potential underserved niches",
                font_size=11, color=MID_GREY)

    gaps = analysis.get("content_gaps", [])
    recommendations = analysis.get("recommendations", [])

    # Left panel: gaps
    add_textbox(slide, Inches(0.4), Inches(1.15), Inches(5.5), Inches(0.4),
                "Topics in Transcripts, Not in Titles", font_size=12, bold=True, color=ACCENT2)
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(0.05), Inches(5.5), ACCENT2)

    if gaps:
        for i, gap in enumerate(gaps[:10]):
            top = Inches(1.65) + i * Inches(0.45)
            add_rect(slide, Inches(0.55), top, Inches(5.2), Inches(0.38),
                     RGBColor(0x14, 0x14, 0x2A))
            add_textbox(slide, Inches(0.65), top + Inches(0.05), Inches(5.0), Inches(0.3),
                        f"• {gap}", font_size=10, color=LIGHT_GREY)
    else:
        add_textbox(slide, Inches(0.55), Inches(1.65), Inches(5.2), Inches(1.0),
                    "No gaps detected this run — transcript coverage may be limited.",
                    font_size=10, color=MID_GREY)

    # Right panel: recommendations
    add_textbox(slide, Inches(6.8), Inches(1.15), Inches(6.2), Inches(0.4),
                "Recommended Video Ideas", font_size=12, bold=True, color=ACCENT)
    add_rect(slide, Inches(6.8), Inches(1.1), Inches(0.05), Inches(5.5), ACCENT)

    if recommendations:
        rec_list = recommendations if isinstance(recommendations, list) else [recommendations]
        for i, rec in enumerate(rec_list[:7]):
            top = Inches(1.65) + i * Inches(0.62)
            add_rect(slide, Inches(6.95), top, Inches(6.05), Inches(0.55),
                     RGBColor(0x1A, 0x1A, 0x35))
            text = rec if isinstance(rec, str) else rec.get("title", str(rec))
            add_textbox(slide, Inches(7.05), top + Inches(0.07), Inches(5.85), Inches(0.45),
                        text, font_size=9, color=LIGHT_GREY)
    else:
        add_textbox(slide, Inches(6.95), Inches(1.65), Inches(6.05), Inches(1.0),
                    "No recommendations written yet.\n"
                    "Complete Phase 2 of the workflow (agent fills in 'recommendations').",
                    font_size=10, color=MID_GREY)

    # Footer
    add_textbox(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
                f"Generated {analysis.get('generated_date', '')} | "
                f"Data period: past {analysis.get('period_days', 14)} days | "
                "Sources: YouTube Data API v3 + youtube-transcript-api",
                font_size=8, color=MID_GREY)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Build YouTube trends .pptx slide deck")
    parser.add_argument("--input", help="Path to analysis JSON (overrides --date)")
    parser.add_argument("--date", help="Date label YYYY-MM-DD (default: most recent)")
    args = parser.parse_args()

    if args.input:
        analysis_path = Path(args.input)
        if not analysis_path.exists():
            print(f"ERROR: {analysis_path} not found", file=sys.stderr)
            sys.exit(1)
    else:
        analysis_path = find_analysis(args.date)

    run_date = analysis_path.stem.replace("analysis_", "") or date.today().isoformat()
    analysis = json.loads(analysis_path.read_text(encoding="utf-8"))

    CHARTS_DIR.mkdir(parents=True, exist_ok=True)
    DECKS_DIR.mkdir(parents=True, exist_ok=True)

    # Create presentation with widescreen (13.33 x 7.5 inches)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("  Building slide 1: Cover...", file=sys.stderr)
    build_slide_cover(prs, analysis)

    print("  Building slide 2: Top Videos...", file=sys.stderr)
    build_slide_top_videos(prs, analysis)

    print("  Building slide 3: View Velocity...", file=sys.stderr)
    build_slide_view_velocity(prs, analysis)

    print("  Building slide 4: Engagement Rate...", file=sys.stderr)
    build_slide_engagement(prs, analysis)

    print("  Building slide 5: Keywords...", file=sys.stderr)
    build_slide_keywords(prs, analysis)

    print("  Building slide 6: Top Channels...", file=sys.stderr)
    build_slide_top_channels(prs, analysis)

    print("  Building slide 7: Transcript Themes...", file=sys.stderr)
    build_slide_transcript_themes(prs, analysis)

    print("  Building slide 8: Content Gaps...", file=sys.stderr)
    build_slide_content_gaps(prs, analysis)

    output_path = DECKS_DIR / f"youtube_trends_{run_date}.pptx"
    prs.save(str(output_path))

    if not output_path.exists() or output_path.stat().st_size == 0:
        print("ERROR: Deck file was not written or is empty.", file=sys.stderr)
        sys.exit(1)

    print(json.dumps({
        "status": "ok",
        "date": run_date,
        "slides": 8,
        "output": str(output_path),
        "file_size_kb": round(output_path.stat().st_size / 1024, 1),
    }, indent=2))


if __name__ == "__main__":
    main()
